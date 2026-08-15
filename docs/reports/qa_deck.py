"""Geometry QA + approximate render, for an environment with no LibreOffice.

Two jobs:
  1. Measure every text frame's wrapped extent against its own box, using real
     font metrics on a CJK font, and report anything that overflows.  Text
     overflow is the defect a rendered preview is normally used to catch.
  2. Draw an approximate raster of each slide from the shape geometry, so
     overlaps, gaps and alignment can be eyeballed.

The raster is an approximation of PowerPoint's layout, not a substitute for
it -- but the overflow numbers come from measuring the same string with the
same point size in the same box, which is the part that actually decides fit.
"""
import sys

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


def qn(tag):
    ns, name = tag.split(":")
    return (A if ns == "a" else P) + name


def solid_of(el):
    """The srgbClr of the nearest solidFill under `el`, or None."""
    if el is None:
        return None
    sf = el.find(".//" + qn("a:solidFill"))
    if sf is None:
        return None
    c = sf.find(qn("a:srgbClr"))
    return c.get("val") if c is not None else None


CJK = "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc"
PX_PER_IN = 110


def font(size_pt, bold=False):
    px = max(int(size_pt * PX_PER_IN / 72), 6)
    return ImageFont.truetype(CJK, px)


def wrap(draw, text, fnt, max_px):
    """Greedy wrap that breaks between CJK chars and on spaces, like a
    renderer does; returns the wrapped lines."""
    lines, cur = [], ""
    for ch in text:
        if ch == "\n":
            lines.append(cur)
            cur = ""
            continue
        trial = cur + ch
        if draw.textlength(trial, font=fnt) > max_px and cur:
            lines.append(cur)
            cur = ch.lstrip() if ch == " " else ch
        else:
            cur = trial
    if cur:
        lines.append(cur)
    return lines


def main(path):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    W, H = int(Emu(sw).inches * PX_PER_IN), int(Emu(sh).inches * PX_PER_IN)
    problems = []

    for idx, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (W, H), "white")
        d = ImageDraw.Draw(img)
        # background: read the solidFill straight out of the XML --
        # python-pptx's fill API does not classify what pptxgenjs writes here
        bgc = solid_of(slide._element.find(qn("p:cSld")))
        if bgc:
            d.rectangle([0, 0, W, H], fill="#" + bgc)

        for shp in slide.shapes:
            if shp.left is None:
                continue
            x0 = Emu(shp.left).inches * PX_PER_IN
            y0 = Emu(shp.top).inches * PX_PER_IN
            w = Emu(shp.width).inches * PX_PER_IN
            h = Emu(shp.height).inches * PX_PER_IN

            fill = solid_of(shp._element.find(qn("p:spPr")))
            fill = "#" + fill if fill else None
            if shp.shape_type is not None and fill:
                if "ELLIPSE" in str(shp.shape_type):
                    d.ellipse([x0, y0, x0 + w, y0 + h], fill=fill)
                else:
                    d.rectangle([x0, y0, x0 + w, y0 + h], fill=fill)

            if shp.has_table:
                tb = shp.table
                ry = y0
                for r in tb.rows:
                    cx = x0
                    for ci, cell in enumerate(r.cells):
                        cw = Emu(tb.columns[ci].width).inches * PX_PER_IN
                        rh = Emu(r.height).inches * PX_PER_IN
                        cf = solid_of(cell._tc.find(qn("a:tcPr")))
                        d.rectangle([cx, ry, cx + cw, ry + rh],
                                    fill="#" + cf if cf else None,
                                    outline="#D3DBE3")
                        para = cell.text_frame.paragraphs[0]
                        runs = para.runs
                        if runs:
                            sz = runs[0].font.size.pt if runs[0].font.size else 12
                            bold = bool(runs[0].font.bold)
                            col = "#000000"
                            rc = solid_of(runs[0]._r.find(qn("a:rPr")))
                            if rc:
                                col = "#" + rc
                            f = font(sz, bold)
                            t = cell.text_frame.text
                            tw = d.textlength(t, font=f)
                            if tw > cw - 6:
                                problems.append(
                                    f"slide{idx}: TABLE cell overflows "
                                    f"({tw:.0f}px in {cw:.0f}px): {t[:40]!r}")
                            d.text((cx + 6, ry + rh / 2), t, font=f, fill=col,
                                   anchor="lm")
                        cx += cw
                    ry += rh
                continue

            if not shp.has_text_frame or not shp.text_frame.text.strip():
                continue

            tf = shp.text_frame
            # pptxgenjs writes margin 0 where we asked; assume a small inset
            inset = 4
            cur_y = y0 + inset
            total_h = 0
            paras = [p for p in tf.paragraphs if p.runs]
            for pi, p in enumerate(paras):
                runs = p.runs
                if not runs:
                    continue
                sz = runs[0].font.size.pt if runs[0].font.size else 14
                bold = bool(runs[0].font.bold)
                col = "#333333"
                rc = solid_of(runs[0]._r.find(qn("a:rPr")))
                if rc:
                    col = "#" + rc
                f = font(sz, bold)
                text = "".join(r.text for r in runs)
                bullet = "•  " if p._pPr is not None and \
                    p._pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar") is not None else ""
                lines = wrap(d, bullet + text, f, w - 2 * inset)
                lh = f.size * 1.28
                for ln in lines:
                    d.text((x0 + inset, cur_y), ln, font=f, fill=col)
                    cur_y += lh
                    total_h += lh
                if pi != len(paras) - 1:     # no spacing after the last
                    total_h += f.size * 0.5
                    cur_y += f.size * 0.5

            if total_h > h + 2:
                problems.append(
                    f"slide{idx}: TEXT overflows box by {total_h - h:.0f}px "
                    f"(needs {total_h:.0f}, box {h:.0f}): "
                    f"{tf.text[:50]!r}")
            if x0 < 0 or y0 < 0 or x0 + w > W + 1 or y0 + h > H + 1:
                problems.append(
                    f"slide{idx}: shape outside slide: {tf.text[:30]!r}")

        img.save(f"qa-{idx}.png")

    print(f"rendered {len(prs.slides)} slides")
    if problems:
        print(f"\n{len(problems)} PROBLEM(S):")
        for p in problems:
            print("  " + p)
    else:
        print("no overflow detected")
    return 1 if problems else 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1] if len(sys.argv) > 1 else "pllsim_report.pptx"))
